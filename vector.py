"""
prepare_styling_rules_final.py

FINAL VERSION – PRODUCTION SAFE

- Uses REAL URLs / files
- Extracts ONLY actionable clothing styling rules
- Enforces strict hygiene rules for vector search
- PRINTS ONLY (no Pinecone, no embeddings yet)
"""

import re
import uuid
import requests
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

# =========================================================
# SOURCES
# =========================================================

STYLING_URLS = [
    "https://www.masarishop.com/newsroom/fashion-style/mix-and-match-outfit-ideas-creative-ways-to-elevate-your-style.html",
    # "https://www.masarishop.com/newsroom/fashion-style/6-clothing-styles-for-you-to-wear.html",
    # "https://www.masarishop.com/newsroom/fashion-style/color-combination-outfit-guide-for-2025-how-to-mix-and-match-perfectly.html",
    # "https://isufashion.com/blogs/news/how-to-mix-combos-to-create-new-outfit-ideas-every-time",
    # "https://withliana.com/how-to-put-together-an-outfit-fashion-101/"
]

FILE_PATHS = [
    # "Tops vs Trousers.docx",
    # "headgears.pdf",
    # "eyewear.pdf",
    # "5. Pants.pptx",
    # "4. Skirt.pptx"
]

# =========================================================
# CONTROLLED VOCABULARY
# =========================================================

FIT_TERMS = ["oversized", "relaxed", "slim", "fitted", "tailored"]
COLOR_TERMS = ["white", "black", "neutral", "beige", "earthy", "brown", "olive"]
STYLE_TERMS = ["classic", "casual", "formal", "athleisure", "minimal", "polished"]
ITEM_TERMS = [
    "t-shirt", "tee", "shirt", "jeans", "denim", "pants", "trousers",
    "skirt", "dress", "jacket", "blazer", "coat", "sweater", "joggers"
]
OCCASION_TERMS = ["casual", "everyday", "office", "formal", "weekend"]
LAYER_TERMS = ["jacket", "blazer", "coat"]
PAIRING_VERBS = ["pair", "combine", "match", "wear", "layer", "add"]

ACCESSORY_TERMS = [
    "necklace", "jewelry", "earrings", "bracelet",
    "choker", "ring", "hat", "scarf", "accessories"
]

# =========================================================
# EXTRACTION HELPERS
# =========================================================

def extract_text_from_url(url):
    try:
        r = requests.get(url, timeout=20)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, "html.parser")

        for t in soup(["script", "style", "nav", "header", "footer", "aside"]):
            t.decompose()

        blocks = soup.find_all(["p", "li", "h2", "h3"])
        return " ".join(b.get_text(" ", strip=True) for b in blocks)
    except Exception as e:
        print(f"[WARN] URL failed: {url} -> {e}")
        return ""

def extract_text_from_docx(path):
    try:
        doc = Document(path)
        return " ".join(p.text for p in doc.paragraphs if p.text.strip())
    except:
        return ""

def extract_text_from_pdf(path):
    try:
        reader = PdfReader(path)
        return " ".join(p.extract_text() or "" for p in reader.pages)
    except:
        return ""

def extract_text_from_pptx(path):
    try:
        prs = Presentation(path)
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    out.append(shape.text)
        return " ".join(out)
    except:
        return ""

# =========================================================
# TEXT PROCESSING
# =========================================================

def split_sentences(text):
    text = re.sub(r"\s+", " ", text.replace("\n", " ")).strip()
    return [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if len(s.strip()) > 12]

def contains_any(terms, text):
    t = text.lower()
    return any(term in t for term in terms)

# =========================================================
# CANONICALIZATION (STRICT)
# =========================================================

def canonicalize_sentence(sentence):
    s = sentence.strip()

    # 1️⃣ REMOVE TITLES / HEADINGS
    if ":" in s:
        left, right = s.split(":", 1)
        if len(left.split()) <= 6:
            s = right.strip()

    s = s.lower()

    # 2️⃣ DROP ACCESSORY-ONLY CONTENT
    if contains_any(ACCESSORY_TERMS, s) and not contains_any(ITEM_TERMS, s):
        return ""

    # normalize wording
    s = re.sub(r'\bbaggy\b', 'oversized', s)
    s = re.sub(r'\bloose\b', 'relaxed', s)
    s = re.sub(r'\bslim-fit\b', 'slim', s)

    s = re.sub(r'\s+', ' ', s).strip()

    # must be actionable
    if contains_any(PAIRING_VERBS, s) and contains_any(ITEM_TERMS, s):
        return s.capitalize() + "."

    # layering rule with clothing
    if contains_any(LAYER_TERMS, s) and contains_any(ITEM_TERMS, s):
        return s.capitalize() + "."

    return ""

# =========================================================
# METADATA EXTRACTION
# =========================================================

def extract_metadata(text):
    t = text.lower()
    meta = {
        "fit": [],
        "color": [],
        "style": [],
        "items": [],
        "occasion": [],
        "layering": False
    }

    for term in FIT_TERMS:
        if term in t:
            meta["fit"].append(term)

    for term in COLOR_TERMS:
        if term in t:
            meta["color"].append(term)

    for term in STYLE_TERMS:
        if term in t:
            meta["style"].append(term)

    for term in ITEM_TERMS:
        if term in t:
            meta["items"].append(term)

    for term in OCCASION_TERMS:
        if term in t:
            meta["occasion"].append(term)

    meta["layering"] = contains_any(LAYER_TERMS, t)

    # implicit inference
    if "blazer" in meta["items"]:
        meta["style"].append("classic")
        meta["layering"] = True

    if "denim" in meta["items"] or "jeans" in meta["items"]:
        meta["style"].append("classic")

    if "white" in meta["color"] or "neutral" in meta["color"]:
        meta["style"].append("minimal")

    if "tailored" in meta["fit"]:
        meta["style"].append("polished")

    for k in meta:
        if isinstance(meta[k], list):
            meta[k] = sorted(set(meta[k]))

    return meta

def compute_confidence(meta):
    score = 0
    if meta["items"]: score += 2
    if meta["fit"]: score += 1
    if meta["color"]: score += 1
    if meta["style"]: score += 1
    if meta["occasion"]: score += 1
    if meta["layering"]: score += 0.5
    return round(min(score / 6, 1), 2)

# =========================================================
# BUILD RECORDS (FINAL GUARDS)
# =========================================================

def build_records(source, text, seen):
    records = []
    for sent in split_sentences(text):

        canonical = canonicalize_sentence(sent)
        if not canonical:
            continue

        meta = extract_metadata(canonical)

        # ❌ HARD RULE: must contain clothing items
        if not meta["items"]:
            continue

        conf = compute_confidence(meta)

        # ❌ HARD RULE: minimum quality
        if conf < 0.45:
            continue

        canon_key = canonical.lower()
        if canon_key in seen:
            continue
        seen.add(canon_key)

        records.append({
            "id": str(uuid.uuid4()),
            "source": source,
            "canonical_text": canonical,
            "metadata": meta,
            "confidence": conf,
            "raw_text": sent
        })

    return records

# =========================================================
# MAIN
# =========================================================

def main():
    all_records = []
    seen_canonicals = set()

    print("\n--- URL SOURCES ---\n")
    for url in STYLING_URLS:
        text = extract_text_from_url(url)
        recs = build_records(url, text, seen_canonicals)
        print(f"{url} -> {len(recs)} rules")
        all_records.extend(recs)

    print(f"\nTOTAL FINAL CLEAN RULES: {len(all_records)}\n")

    for r in all_records:
        print("--------------------------------------------------------")
        print("ID:", r["id"])
        print("SOURCE:", r["source"])
        print("CONFIDENCE:", r["confidence"])
        print("CANONICAL:", r["canonical_text"])
        print("METADATA:", r["metadata"])
        print("RAW:", r["raw_text"])

if __name__ == "__main__":
    main()
